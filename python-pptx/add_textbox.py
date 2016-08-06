from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

# normal textbox
tf.text = "This is text inside a textbox"

# textbox with bold text
p = tf.add_paragraph()
p.text = "This is a second paragraph that's bold"
p.font.bold = True

# textbox with Big size text
p = tf.add_paragraph()
p.text = "This is a third paragraph that's big"
p.font.size = Pt(40)

# textbox has none margin
left = top = width = height = Inches(0.05)
txBox2 = slide.shapes.add_textbox(left, top, width, height)
txBox2.left = 2
txBox2.top = 2
tf = txBox2.text_frame
p.text = "This is a fourth paragraph that's no margin"

prs.save('test_with_textbox.pptx')